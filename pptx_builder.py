"""
CHI Insurance Brokers — PPTX Slide Builder

Βελτιώσεις σε σχέση με v1:
  • Διαχωρισμός σε επιμέρους συναρτήσεις (cover, overview, proposal, table, closing)
  • Μηνιαίο ασφάλιστρο εμφανίζεται σε κάρτες και per-proposal slides
  • Υποστήριξη συχνότητας πληρωμής (Μηνιαία / Τριμηνιαία / Εξαμηνιαία / Ετήσια)
  • Βαθμολογία κάλυψης (0–10) σε κάθε κάρτα και slide προσφοράς
  • Μοναδική κάλυψη (★) στον πίνακα σύγκρισης
  • 14 γραμμές στον πίνακα με νέα πεδία
  • Παρατηρήσεις στο closing slide
  • Υποστήριξη λογότυπου
  • Διόρθωση null guard για outpatient_pct
  • Σωστός τίτλος για 2 προτάσεις
  • Πλήρως ελληνική ορολογία
"""

import io
from datetime import datetime

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from config import C, insurer_color, rgb
from extraction import compute_score


# ─── LOW-LEVEL DRAWING HELPERS ──────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    """Add a filled rectangle at (x, y) inches with size (w, h) inches."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb  = line_color
        shape.line.width      = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=12, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, valign="middle", wrap=True):
    """Add a text box at (x, y) inches."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap

    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM

    p          = tf.paragraphs[0]
    p.alignment = align
    run        = p.add_run()
    run.text   = str(text)
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


# ─── SHARED SLIDE UTILITIES ─────────────────────────────────────────

def _footer(slide, text, light=False):
    add_rect(slide, 0, 7.18, 13.33, 0.32, C["navy"] if light else C["navyDark"])
    add_text(slide, text, 0.3, 7.18, 12.7, 0.32,
             size=8, color=C["teal"], align=PP_ALIGN.CENTER)


def _top_bar(slide, color):
    add_rect(slide, 0, 0, 13.33, 0.1, color)


def _add_logo(slide, logo_bytes, x, y, h_in):
    """Safely place a logo image; silently skip if bytes are invalid."""
    if not logo_bytes:
        return
    try:
        slide.shapes.add_picture(
            io.BytesIO(logo_bytes), Inches(x), Inches(y), height=Inches(h_in)
        )
    except Exception:
        pass


def _sym(prop: dict) -> str:
    cur = prop.get("currency") or "EUR"
    return "€" if cur == "EUR" else ("$" if cur == "USD" else "£")


def _monthly(prop: dict):
    """Return monthly equivalent of annual_premium, or None."""
    try:
        annual = int(str(prop.get("annual_premium", 0) or 0).replace(",", ""))
        return round(annual / 12) if annual else None
    except (ValueError, TypeError):
        return None


# Χαρτογράφηση συχνότητας πληρωμής → παρονομαστής ανά έτος
_FREQ_DIVISOR = {
    "Μηνιαία":      12,
    "Τριμηνιαία":    4,
    "Εξαμηνιαία":   2,
    "Ετήσια":        1,
}
_FREQ_LABEL = {
    "Μηνιαία":     "μηνιαίο ασφάλιστρο",
    "Τριμηνιαία":  "τριμηνιαίο ασφάλιστρο",
    "Εξαμηνιαία":  "εξαμηνιαίο ασφάλιστρο",
    "Ετήσια":       "ετήσιο ασφάλιστρο",
}


def _payment_display(prop: dict):
    """
    Return (display_amount_str, frequency_label, monthly_str_or_None).

    Βασίζεται στο annual_premium (πάντα ετήσιο) και στο payment_frequency
    που επέλεξε ο χρήστης.
    """
    sym  = _sym(prop)
    freq = prop.get("payment_frequency") or "Ετήσια"
    try:
        annual   = int(str(prop.get("annual_premium", 0) or 0).replace(",", ""))
        divisor  = _FREQ_DIVISOR.get(freq, 1)
        per_freq = round(annual / divisor) if annual else None
        label    = _FREQ_LABEL.get(freq, "ετήσιο ασφάλιστρο")

        display  = f"{sym}{per_freq:,}" if per_freq else "—"

        # Δευτερεύουσα ένδειξη: αν δεν εμφανίζεται ήδη μηνιαία, δείξε μηνιαία
        if freq != "Μηνιαία" and annual:
            monthly_val = round(annual / 12)
            monthly_str = f"{sym}{monthly_val:,} / μήνα"
        else:
            monthly_str = None

        return display, label, monthly_str
    except (ValueError, TypeError):
        raw = prop.get("annual_premium", "—")
        return f"{sym}{raw}", "ετήσιο ασφάλιστρο", None


def _is_covered(val) -> bool:
    if val is None:
        return False
    s = str(val).strip()
    return s not in ("", "null", "None", "—") and "Not Covered" not in s


def _display(val) -> str:
    return str(val) if _is_covered(val) else "ΔΕΝ καλύπτει"


def _n_props_label(n: int) -> str:
    return {2: "Δύο", 3: "Τρεις", 4: "Τέσσερις"}.get(n, str(n))


def _fmt_outpatient(prop: dict, sym: str) -> str:
    """Format outpatient coverage string, guarding against null % value."""
    ol  = prop.get("outpatient_limit")
    pct = prop.get("outpatient_pct")
    if not ol or str(ol).strip() in ("null", "None", "", "Not Covered"):
        return ol or "Not Covered"
    pct_str = (
        f" ({pct}%)"
        if pct and str(pct).strip() not in ("null", "None", "")
        else ""
    )
    return f"εως {sym}{ol}{pct_str}"


# ─── SLIDE 1: COVER ─────────────────────────────────────────────────

def _build_cover(prs, layout, client_name, client_members, proposals,
                 footer_text, logo_bytes):
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navyDark"]

    _top_bar(s, C["teal"])
    add_rect(s, 0, 0.1, 5.0, 7.08, C["navy"])
    add_rect(s, 5.0, 0.1, 0.05, 7.08, C["teal"])

    # Optional logo on left panel
    if logo_bytes:
        _add_logo(s, logo_bytes, 0.5, 0.9, 1.1)

    logo_offset = 1.25 if logo_bytes else 0.0
    add_text(s, "ΑΣΦΑΛΙΣΗ ΥΓΕΙΑΣ", 0.3, 2.6 + logo_offset * 0.3, 4.4, 0.5,
             size=12, bold=True, color=C["teal"], align=PP_ALIGN.CENTER)
    add_text(s, "ΣΥΓΚΡΙΤΙΚΗ ΑΝΑΛΥΣΗ", 0.3, 3.1 + logo_offset * 0.3, 4.4, 0.4,
             size=10, color=rgb(0x9D, 0xC4, 0xD8), align=PP_ALIGN.CENTER)

    # Right panel — client details
    add_text(s, "Πρόταση", 5.5, 1.0, 7.5, 0.8,
             size=22, italic=True, color=C["teal"])
    add_text(s, "Ασφάλισης Υγείας", 5.5, 1.7, 7.5, 1.2,
             size=44, bold=True, color=C["white"])
    add_rect(s, 5.5, 3.0, 7.6, 0.05, C["teal"])

    add_text(s, client_name, 5.5, 3.15, 7.5, 0.6,
             size=20, bold=True, color=C["white"])
    member_str = "  ·  ".join(
        f"{m.get('role', '')} ({m.get('age', '')} ετών)" for m in client_members
    )
    add_text(s, member_str, 5.5, 3.75, 7.5, 0.45,
             size=12, color=rgb(0x9D, 0xC4, 0xD8))

    # Date badge
    add_rect(s, 5.5, 4.4, 2.1, 0.45, C["teal"])
    add_text(s, datetime.now().strftime("%B %Y"), 5.5, 4.4, 2.1, 0.45,
             size=11, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)

    # Insurer name chips along bottom
    for i, prop in enumerate(proposals[:4]):
        xp  = 5.5 + i * 1.95
        col = insurer_color(prop.get("insurer", ""))
        add_rect(s, xp, 5.2, 1.82, 0.42, col)
        add_text(s, prop.get("insurer", "").upper(), xp, 5.2, 1.82, 0.42,
                 size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    _footer(s, footer_text, light=True)


# ─── SLIDE 2: OVERVIEW CARDS ────────────────────────────────────────

def _build_overview(prs, layout, client_name, proposals, recommended_idx, footer_text):
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navy"]
    _top_bar(s, C["gold"])

    n = len(proposals)
    add_text(s, f"{_n_props_label(n)} Προτάσεις για {client_name}",
             0.4, 0.18, 12.5, 0.62, size=28, bold=True, color=C["white"])
    add_text(s, "Από την οικονομική επιλογή έως την ολοκληρωμένη διεθνή κάλυψη",
             0.4, 0.82, 12.5, 0.38, size=12, italic=True,
             color=rgb(0x9D, 0xC4, 0xD8))

    card_w = (13.33 - 0.8 - (n - 1) * 0.25) / n

    for i, prop in enumerate(proposals):
        xp     = 0.4 + i * (card_w + 0.25)
        is_rec = (i == recommended_idx)
        col    = insurer_color(prop.get("insurer", ""))
        bg_col = rgb(0x1A, 0x2F, 0x45) if is_rec else rgb(0x17, 0x35, 0x4E)
        border = C["gold"] if is_rec else col
        score  = compute_score(prop)
        sym    = _sym(prop)
        display_amt, freq_lbl, monthly_str = _payment_display(prop)

        # Card body
        add_rect(s, xp, 1.38, card_w, 5.62, bg_col,
                 line_color=border, line_width=2.0 if is_rec else 0.8)
        # Header bar (insurer color)
        add_rect(s, xp, 1.38, card_w, 0.5, col)

        label = "★ ΠΡΟΤΕΙΝΟΜΕΝΗ" if is_rec else f"ΕΠΙΛΟΓΗ {chr(65 + i)}"
        add_text(s, label, xp, 1.38, card_w, 0.5,
                 size=9, bold=True,
                 color=C["navy"] if is_rec else C["white"],
                 align=PP_ALIGN.CENTER)

        add_text(s, prop.get("insurer", ""), xp + 0.1, 1.93, card_w - 0.2, 0.38,
                 size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, prop.get("plan_name", ""), xp + 0.1, 2.31, card_w - 0.2, 0.58,
                 size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        # Ασφάλιστρο (κύριο)
        price_bg = C["gold"] if is_rec else rgb(0x2A, 0x4A, 0x63)
        price_fg = C["navy"] if is_rec else C["teal"]
        add_rect(s, xp + card_w * 0.1, 2.98, card_w * 0.8, 0.56, price_bg)
        add_text(s, display_amt,
                 xp + card_w * 0.1, 2.98, card_w * 0.8, 0.56,
                 size=20, bold=True, color=price_fg, align=PP_ALIGN.CENTER)

        # Συχνότητα / μηνιαία ισοδύναμο
        add_text(s, freq_lbl,
                 xp + 0.1, 3.57, card_w - 0.2, 0.24,
                 size=7, italic=True,
                 color=rgb(0x9D, 0xC4, 0xD8), align=PP_ALIGN.CENTER)
        if monthly_str:
            add_text(s, monthly_str,
                     xp + 0.1, 3.82, card_w - 0.2, 0.24,
                     size=7, italic=True,
                     color=rgb(0x9D, 0xC4, 0xD8), align=PP_ALIGN.CENTER)

        # Βαθμός κάλυψης
        sc_col = C["green"] if score >= 7 else (C["orange"] if score >= 5 else C["red"])
        add_rect(s, xp + card_w * 0.2, 4.1, card_w * 0.6, 0.28, sc_col)
        add_text(s, f"Βαθμός: {score} / 10",
                 xp + card_w * 0.2, 4.1, card_w * 0.6, 0.28,
                 size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        # Γεωγραφία
        add_text(s, prop.get("geography", ""),
                 xp + 0.1, 4.42, card_w - 0.2, 0.3,
                 size=8, italic=True,
                 color=rgb(0x9D, 0xC4, 0xD8), align=PP_ALIGN.CENTER)

        # Συνοπτικά bullets
        bullets = [
            f"Απαλλαγή: {prop.get('deductible') or '—'}",
            f"Νοσηλεία: {prop.get('inpatient') or '—'}",
            f"MRI/PET:  {prop.get('mri_ct_pet') or '—'}",
            f"Αναμονή:  {prop.get('waiting_period') or '—'}",
        ]
        for bi, b in enumerate(bullets):
            add_text(s, f"• {b}", xp + 0.15, 4.76 + bi * 0.44, card_w - 0.3, 0.4,
                     size=8, color=rgb(0xC8, 0xDF, 0xF0))

    _footer(s, footer_text, light=True)


# ─── SLIDES 3+: ONE PER PROPOSAL ────────────────────────────────────

def _build_proposal_slide(prs, layout, prop, idx, is_rec, footer_text):
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navyDark"] if is_rec else C["offWhite"]

    col = insurer_color(prop.get("insurer", ""))
    _top_bar(s, C["gold"] if is_rec else col)
    add_rect(s, 0, 0.1, 0.35, 7.08, C["gold"] if is_rec else col)

    label = (
        f"★ ΕΠΙΛΟΓΗ {chr(65 + idx)} — ΠΡΟΤΕΙΝΟΜΕΝΗ"
        if is_rec else f"ΕΠΙΛΟΓΗ {chr(65 + idx)}"
    )
    add_text(s, label, 0.5, 0.15, 8.0, 0.42,
             size=10, bold=True, color=C["gold"] if is_rec else col)
    add_text(s, f"{prop.get('insurer', '')} — {prop.get('plan_name', '')}",
             0.5, 0.55, 9.5, 0.7, size=23, bold=True,
             color=C["white"] if is_rec else C["navy"])

    # Πλαίσιο τιμής (πάνω-δεξιά)
    sym     = _sym(prop)
    display_amt, freq_lbl, monthly_str = _payment_display(prop)
    price_col = C["gold"] if is_rec else col
    add_rect(s, 10.15, 0.25, 2.85, 1.1, price_col)
    add_text(s, display_amt,
             10.15, 0.25, 2.85, 0.68,
             size=25, bold=True,
             color=C["navy"] if is_rec else C["white"],
             align=PP_ALIGN.CENTER)
    add_text(s, freq_lbl,
             10.15, 0.91, 2.85, 0.28, size=8,
             color=rgb(0xC8, 0xDF, 0xF0) if is_rec else rgb(0x5A, 0x4A, 0x00),
             align=PP_ALIGN.CENTER)

    # Μηνιαίο ισοδύναμο (αν διαφέρει)
    if monthly_str:
        add_text(s, monthly_str,
                 10.15, 1.37, 2.85, 0.26, size=8, italic=True,
                 color=rgb(0xC8, 0xDF, 0xF0) if is_rec else rgb(0x5A, 0x4A, 0x00),
                 align=PP_ALIGN.CENTER)

    # Βαθμός κάλυψης
    score  = compute_score(prop)
    sc_col = C["green"] if score >= 7 else (C["orange"] if score >= 5 else C["red"])
    add_rect(s, 10.55, 1.68, 2.1, 0.32, sc_col)
    add_text(s, f"Βαθμός: {score} / 10",
             10.55, 1.68, 2.1, 0.32,
             size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    # Key parameters (top-left grid)
    params = [
        ("Μέγιστο Κεφάλαιο", f"{sym}{prop.get('max_coverage') or '—'}"),
        ("Απαλλαγή",          prop.get("deductible") or "—"),
        ("Γεωγραφία",         prop.get("geography") or "—"),
        ("Θέση Νοσηλείας",    prop.get("hospital_class") or "—"),
    ]
    for pi, (k, v) in enumerate(params):
        yp = 1.55 + pi * 0.5
        lbl_bg = C["navy"] if not is_rec else rgb(0x1A, 0x35, 0x50)
        lbl_fg = C["white"] if not is_rec else C["gold"]
        add_rect(s, 0.5, yp, 3.1, 0.42, lbl_bg)
        add_text(s, k, 0.5, yp, 3.1, 0.42,
                 size=9, bold=True, color=lbl_fg, align=PP_ALIGN.CENTER)
        add_text(s, v, 3.7, yp + 0.04, 6.2, 0.42,
                 size=9,
                 color=C["textDark"] if not is_rec else rgb(0xC8, 0xDF, 0xF0))

    # ── LEFT PANEL: Coverages ──
    add_rect(s, 0.5, 3.63, 6.0, 0.36, C["green"])
    add_text(s, "✓  ΚΑΛΥΨΕΙΣ", 0.5, 3.63, 6.0, 0.36,
             size=9, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    covers = [
        ("Νοσηλεία",           prop.get("inpatient") or "—"),
        ("Εξωνοσοκ.",          _fmt_outpatient(prop, sym)),
        ("MRI / CT / PET",     prop.get("mri_ct_pet") or "—"),
        ("Καρκίνος",            prop.get("cancer") or "—"),
        ("Χρόνιες Παθήσεις",   prop.get("chronic_conditions") or "—"),
        ("Εκκένωση/Μεταφορά",  prop.get("evacuation_repatriation") or "—"),
        ("Φυσιοθεραπεία",       prop.get("physiotherapy") or "—"),
    ]
    for ci, (k, v) in enumerate(covers):
        yp   = 4.03 + ci * 0.41
        tick = "✓" if _is_covered(v) else "✗"
        tc   = C["green"] if tick == "✓" else C["red"]
        add_text(s, tick, 0.55, yp, 0.4, 0.37,
                 size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
        add_text(s, k, 1.0, yp, 1.85, 0.37,
                 size=9, bold=True,
                 color=C["white"] if is_rec else C["navy"])
        add_text(s, v, 2.9, yp, 3.5, 0.37,
                 size=8,
                 color=rgb(0xC8, 0xDF, 0xF0) if is_rec else C["textDark"])

    # ── RIGHT PANEL: Notes ──
    add_rect(s, 6.8, 3.63, 6.2, 0.36, C["orange"])
    add_text(s, "⚠  ΣΗΜΑΝΤΙΚΕΣ ΠΑΡΑΤΗΡΗΣΕΙΣ", 6.8, 3.63, 6.2, 0.36,
             size=9, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    extra_notes = [
        f"Ψυχ. Νοσηλεία: {prop.get('psychiatric_inpatient') or 'Δ/Α'}",
        f"Οδοντιατρική Έκτ.: {prop.get('dental_emergency') or 'Δ/Α'}",
        f"Αναμονή: {prop.get('waiting_period') or 'Άμεση'}",
        f"Προϋπ. παθήσεις: {prop.get('preexisting') or '—'}",
    ]
    all_notes = list(prop.get("key_notes") or []) + extra_notes

    for ni, note in enumerate(all_notes[:7]):
        yp = 4.03 + ni * 0.41
        add_text(s, f"• {note}", 6.85, yp, 6.1, 0.38,
                 size=8.5,
                 color=rgb(0xC8, 0xDF, 0xF0) if is_rec else C["textDark"])

    _footer(s, footer_text, light=is_rec)


# ─── COMPARISON TABLE ───────────────────────────────────────────────

def _build_comparison_table(prs, layout, proposals, recommended_idx, footer_text):
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["offWhite"]
    _top_bar(s, C["teal"])

    add_text(s, "Σύγκριση Καλύψεων — Πίνακας",
             0.4, 0.12, 12.5, 0.55, size=23, bold=True, color=C["navy"])
    add_text(s, "★ = Μοναδική κάλυψη σε αυτό το πλάνο  |  Έντονο = Προτεινόμενη επιλογή",
             0.4, 0.67, 12.5, 0.28, size=8, italic=True, color=C["orange"])

    n = len(proposals)
    label_w  = 3.2
    data_w   = (13.33 - 0.8 - label_w) / n
    col_widths = [label_w] + [data_w] * n
    xstarts    = [0.4]
    for cw in col_widths[:-1]:
        xstarts.append(xstarts[-1] + cw)

    # Column headers
    header_colors      = [C["navy"]] + [insurer_color(p.get("insurer", "")) for p in proposals]
    header_text_colors = [C["white"]] * len(col_widths)
    if recommended_idx < n:
        header_colors[recommended_idx + 1]      = C["gold"]
        header_text_colors[recommended_idx + 1] = C["navy"]

    headers = ["ΚΑΛΥΨΗ"] + [
        f"{p.get('insurer','')}\n{p.get('plan_name','')}" for p in proposals
    ]
    for h, xp, cw, hc, htc in zip(
        headers, xstarts, col_widths, header_colors, header_text_colors
    ):
        add_rect(s, xp, 1.0, cw, 0.58, hc)
        add_text(s, h, xp + 0.05, 1.0, cw - 0.1, 0.58,
                 size=7.5, bold=True, color=htc, align=PP_ALIGN.CENTER)

    # Data rows — 14 rows, 0.37 height each → fits on slide
    row_h   = 0.37
    start_y = 1.62

    row_data = [
        ("Μέγ. Κεφάλαιο / Έτος",  lambda p: f"{_sym(p)}{p.get('max_coverage','—') or '—'}"),
        ("Απαλλαγή",               lambda p: p.get("deductible", "—") or "—"),
        ("Γεωγραφία",              lambda p: p.get("geography", "—") or "—"),
        ("Θέση Νοσηλείας",         lambda p: p.get("hospital_class", "—") or "—"),
        ("Νοσηλεία Εσωτερικού",    lambda p: p.get("inpatient", "—") or "—"),
        ("Χρόνιες Παθήσεις",       lambda p: p.get("chronic_conditions", "—") or "—"),
        ("MRI / CT / PET",         lambda p: p.get("mri_ct_pet", "—") or "—"),
        ("Εξωνοσ. Καλύψεις",      lambda p: _fmt_outpatient(p, _sym(p))),
        ("Φυσιοθεραπεία",          lambda p: p.get("physiotherapy", "—") or "—"),
        ("Καρκίνος",                lambda p: p.get("cancer", "—") or "—"),
        ("Ψυχ. Νοσηλεία",          lambda p: p.get("psychiatric_inpatient", "—") or "—"),
        ("Οδοντιατρική Έκτακτη",   lambda p: p.get("dental_emergency", "—") or "—"),
        ("Εκκένωση / Μεταφορά",   lambda p: p.get("evacuation_repatriation", "—") or "—"),
        ("Ετήσιο Ασφάλιστρο",     lambda p: f"{_sym(p)}{p.get('annual_premium','—') or '—'}"),
    ]

    for ri, (label, fn) in enumerate(row_data):
        yp = start_y + ri * row_h

        # Alternating row backgrounds
        row_bg = rgb(0xF4, 0xF9, 0xFF) if ri % 2 == 0 else C["white"]

        # Compute values and detect unique coverage per proposal
        vals          = [fn(p) for p in proposals]
        covered_flags = [_is_covered(v) for v in vals]
        n_covered     = sum(covered_flags)
        unique        = [covered_flags[pi] and n_covered == 1 for pi in range(n)]

        for ci, (xp, cw) in enumerate(zip(xstarts, col_widths)):
            is_rec_col = (ci == recommended_idx + 1)
            bg = (
                (rgb(0xFF, 0xF8, 0xE1) if ri % 2 == 0 else rgb(0xFF, 0xF3, 0xC4))
                if is_rec_col else row_bg
            )
            add_rect(s, xp, yp, cw, row_h, bg)

            if ci == 0:
                # Row label
                add_text(s, label, xp + 0.1, yp, cw - 0.12, row_h,
                         size=8.5, bold=True, color=C["navy"])
            else:
                pi      = ci - 1
                val     = vals[pi]
                display = _display(val)
                is_good = _is_covered(val)

                # Unique-coverage teal accent bar on the left edge
                if unique[pi]:
                    add_rect(s, xp, yp, 0.05, row_h, C["teal"])
                    display = f"★ {display}"

                # Text color
                if any(c in display for c in ("€", "$", "£")):
                    text_color = C["navy"]
                elif is_good:
                    text_color = C["green"]
                else:
                    text_color = C["red"]

                add_text(s, display, xp + 0.06, yp, cw - 0.1, row_h,
                         size=7.5, bold=is_rec_col,
                         color=text_color, align=PP_ALIGN.CENTER)

    _footer(s, footer_text)


# ─── CLOSING SLIDE ──────────────────────────────────────────────────

def _build_closing(prs, layout, proposals, recommended_idx, footer_text):
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navyDark"]
    _top_bar(s, C["gold"])

    add_text(s, "Η Πρότασή μας", 0.5, 0.2, 12.0, 0.65,
             size=30, bold=True, color=C["gold"])

    steps_y = 2.0   # default y-position for the 3 steps

    if recommended_idx < len(proposals):
        rec     = proposals[recommended_idx]
        sym     = _sym(rec)
        display_amt, freq_lbl, monthly_str = _payment_display(rec)
        monthly_part = f"  ·  {monthly_str}" if monthly_str else ""
        rec_text = (
            f"{rec.get('insurer', '')} — {rec.get('plan_name', '')}  "
            f"|  {display_amt} / {freq_lbl}{monthly_part}"
        )
        add_text(s, rec_text, 0.5, 0.92, 12.5, 0.55,
                 size=15, bold=True, color=C["white"])

        # Key advantages from recommended plan's key_notes
        key_notes = [n for n in (rec.get("key_notes") or []) if n]
        if key_notes:
            add_rect(s, 0.5, 1.58, 12.33, 0.32, rgb(0x1A, 0x35, 0x50))
            add_rect(s, 0.5, 1.58, 0.06, 0.32, C["gold"])
            add_text(s, "🔑  ΚΥΡΙΑ ΠΛΕΟΝΕΚΤΗΜΑΤΑ", 0.6, 1.58, 12.0, 0.32,
                     size=9, bold=True, color=C["gold"])

            for ni, note in enumerate(key_notes[:3]):
                add_text(s, f"✓  {note}", 0.65, 1.95 + ni * 0.37, 12.0, 0.34,
                         size=10, color=rgb(0xC8, 0xDF, 0xF0))

            steps_y = 3.15   # push steps down to make room for notes

    # Next steps
    steps = [
        ("ΒΗΜΑ 1", "Εντός 48ωρών",  "Έγκριση πρότασης & αποστολή ιατρικού ιστορικού"),
        ("ΒΗΜΑ 2", "Underwriting",   "Υπογραφή αίτησης — σαφής γνώση τι καλύπτεται"),
        ("ΒΗΜΑ 3", "Ενεργοποίηση",  "Άμεση κάλυψη χωρίς αναμονές"),
    ]
    for si, (tag, title, body) in enumerate(steps):
        xp = 0.5 + si * 4.25
        add_rect(s, xp, steps_y, 4.0, 2.5, rgb(0x1A, 0x35, 0x50))
        add_rect(s, xp, steps_y, 4.0, 0.56, C["gold"])
        add_text(s, tag, xp, steps_y, 4.0, 0.56,
                 size=12, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_text(s, title, xp + 0.1, steps_y + 0.62, 3.8, 0.5,
                 size=13, bold=True, color=C["teal"], align=PP_ALIGN.CENTER)
        add_text(s, body, xp + 0.15, steps_y + 1.18, 3.7, 1.2,
                 size=10, color=rgb(0xC8, 0xDF, 0xF0), align=PP_ALIGN.CENTER)

    # Inspirational quote
    quote_y = steps_y + 2.62
    add_rect(s, 0.5, quote_y, 12.33, 1.0, rgb(0x1A, 0x35, 0x50))
    add_rect(s, 0.5, quote_y, 0.06, 1.0, C["gold"])
    add_text(
        s,
        "Η ασφάλεια υγείας δεν είναι κόστος — "
        "είναι επένδυση στην ηρεμία σας και στην οικογένειά σας.",
        0.7, quote_y, 12.0, 1.0, size=13, italic=True,
        color=rgb(0xC8, 0xDF, 0xF0),
    )

    _footer(s, footer_text, light=True)


# ─── PUBLIC ENTRY POINT ─────────────────────────────────────────────

def _build_analysis_slide(prs, layout, analysis: dict, proposals: list,
                          recommended_idx: int, footer_text: str):
    """
    Dedicated slide: 'Γιατί Προτείνουμε Αυτό το Πλάνο'
    Shows the Claude-generated narrative, key reasons, plan verdicts, and concerns.
    """
    s = prs.slides.add_slide(layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["offWhite"]
    _top_bar(s, C["teal"])

    # ── Title bar ──
    add_rect(s, 0, 0.1, 13.33, 0.75, C["navy"])
    add_text(s, "🔍  Γιατί Προτείνουμε Αυτό το Πλάνο",
             0.3, 0.1, 10.0, 0.75,
             size=20, bold=True, color=C["white"])

    rec = proposals[recommended_idx] if recommended_idx < len(proposals) else {}
    rec_label = f"{rec.get('insurer','?')} — {rec.get('plan_name','?')}"
    add_text(s, rec_label, 10.0, 0.1, 3.0, 0.75,
             size=13, bold=True, color=C["gold"], align=PP_ALIGN.RIGHT)

    # ── Headline banner ──
    headline = analysis.get("headline", "")
    if headline:
        add_rect(s, 0.3, 0.95, 12.73, 0.45, C["teal"])
        add_text(s, headline, 0.45, 0.95, 12.5, 0.45,
                 size=11, bold=True, color=C["white"], italic=True)

    # ── LEFT COLUMN: rationale + key reasons ──
    lx, ly = 0.3, 1.52

    # Main rationale (trimmed to 3 sentences for slide brevity)
    rationale = analysis.get("main_rationale", "")
    sentences = [s2.strip() for s2 in rationale.replace("\n", " ").split(".") if s2.strip()]
    short_rationale = ". ".join(sentences[:3]) + ("." if sentences[:3] else "")

    add_rect(s, lx, ly, 6.2, 0.3, C["navy"])
    add_text(s, "  ΑΙΤΙΟΛΟΓΗΣΗ ΠΡΟΤΑΣΗΣ", lx, ly, 6.2, 0.3,
             size=9, bold=True, color=C["teal"])
    add_rect(s, lx, ly + 0.3, 6.2, 1.5, rgb(0xE8, 0xF4, 0xFF))
    add_text(s, short_rationale, lx + 0.12, ly + 0.32, 5.95, 1.45,
             size=9, color=C["textDark"], italic=True)

    # Key reasons
    reasons = analysis.get("key_reasons", [])[:4]
    ry = ly + 1.92
    add_rect(s, lx, ry, 6.2, 0.3, C["navy"])
    add_text(s, "  ΒΑΣΙΚΟΙ ΛΟΓΟΙ ΕΠΙΛΟΓΗΣ", lx, ry, 6.2, 0.3,
             size=9, bold=True, color=C["gold"])
    for ni, reason in enumerate(reasons):
        row_y = ry + 0.3 + ni * 0.55
        bg = rgb(0xF0, 0xF7, 0xFF) if ni % 2 == 0 else rgb(0xE3, 0xF0, 0xFA)
        add_rect(s, lx, row_y, 6.2, 0.52, bg)
        add_rect(s, lx, row_y, 0.05, 0.52, C["teal"])
        add_text(s, f"✓  {reason}", lx + 0.12, row_y + 0.02, 6.0, 0.48,
                 size=8.5, color=C["textDark"])

    # ── RIGHT COLUMN: plan verdicts + concerns ──
    rx, ry2 = 6.83, 1.52
    tag_colors = {
        "ΑΡΙΣΤΟ":       C["green"],
        "ΚΑΛΟ":         C["teal"],
        "ΜΕΣΑΙΟ":       C["orange"],
        "ΠΕΡΙΟΡΙΣΜΕΝΟ": C["red"],
    }

    verdicts = analysis.get("plan_verdicts", [])
    add_rect(s, rx, ry2, 6.2, 0.3, C["navy"])
    add_text(s, "  ΑΞΙΟΛΟΓΗΣΗ ΠΡΟΣΦΟΡΩΝ", rx, ry2, 6.2, 0.3,
             size=9, bold=True, color=C["teal"])

    vstart = ry2 + 0.3
    row_h  = 0.65
    for vi, v in enumerate(verdicts[:4]):
        vrow_y = vstart + vi * row_h
        color  = tag_colors.get(v.get("tag", ""), C["teal"])
        add_rect(s, rx, vrow_y, 6.2, row_h - 0.04, rgb(0xF4, 0xF9, 0xFF))
        add_rect(s, rx, vrow_y, 0.07, row_h - 0.04, color)
        plan_title = f"{v.get('insurer','')} — {v.get('plan','')}"
        add_text(s, plan_title, rx + 0.15, vrow_y + 0.02, 4.5, 0.26,
                 size=8.5, bold=True, color=C["navy"])
        tag_label = v.get("tag", "")
        add_rect(s, rx + 4.75, vrow_y + 0.04, 1.35, 0.22, color)
        add_text(s, tag_label, rx + 4.75, vrow_y + 0.04, 1.35, 0.22,
                 size=7.5, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        verdict_text = v.get("verdict", "")
        add_text(s, verdict_text, rx + 0.15, vrow_y + 0.3, 5.95, 0.3,
                 size=8, color=rgb(0x44, 0x44, 0x44))

    # Concerns at bottom right
    concerns = analysis.get("key_concerns", [])
    if concerns:
        c_y = vstart + len(verdicts[:4]) * row_h + 0.15
        if c_y < 6.4:
            add_rect(s, rx, c_y, 6.2, 0.28, rgb(0xFF, 0xF3, 0xCD))
            add_rect(s, rx, c_y, 0.07, 0.28, C["orange"])
            add_text(s, "  ΣΗΜΕΙΑ ΠΡΟΣΟΧΗΣ", rx, c_y, 6.2, 0.28,
                     size=8.5, bold=True, color=C["orange"])
            for ci, concern in enumerate(concerns[:2]):
                c_item_y = c_y + 0.28 + ci * 0.38
                if c_item_y < 6.8:
                    add_rect(s, rx, c_item_y, 6.2, 0.35, rgb(0xFF, 0xF8, 0xE8))
                    add_text(s, f"⚠  {concern}", rx + 0.12, c_item_y + 0.02,
                             5.95, 0.31, size=8, color=rgb(0x80, 0x40, 0x00))

    # Decision factors strip at bottom
    factors = analysis.get("decision_factors", [])
    if factors:
        add_rect(s, 0.3, 6.78, 12.73, 0.3, C["navy"])
        factors_str = "  ·  ".join(factors[:3])
        add_text(s, f"Κριτήρια: {factors_str}",
                 0.5, 6.78, 12.5, 0.3,
                 size=8, color=C["teal"])

    _footer(s, footer_text)


def generate_pptx(
    client_name: str,
    client_members: list,
    proposals: list,
    recommended_idx: int,
    broker_name: str,
    broker_tel: str,
    broker_email: str,
    logo_bytes: bytes = None,
    analysis: dict = None,
) -> bytes:
    """
    Build a multi-slide PPTX comparison presentation and return raw bytes.

    Slides:
        1  Cover
        2  Overview cards (all proposals side-by-side)
        3+ One detail slide per proposal
        N-2 Comparison table (all fields)
        N-1 Analysis & Recommendation Rationale  ← NEW (if analysis provided)
        N   Closing / call-to-action
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank  = prs.slide_layouts[6]
    footer = (
        f"CHI Insurance Brokers | {broker_name}  ·  {broker_tel}  ·  {broker_email}"
    )

    _build_cover(prs, blank, client_name, client_members, proposals, footer, logo_bytes)
    _build_overview(prs, blank, client_name, proposals, recommended_idx, footer)

    for i, prop in enumerate(proposals):
        _build_proposal_slide(prs, blank, prop, i, i == recommended_idx, footer)

    _build_comparison_table(prs, blank, proposals, recommended_idx, footer)

    if analysis:
        _build_analysis_slide(prs, blank, analysis, proposals, recommended_idx, footer)

    _build_closing(prs, blank, proposals, recommended_idx, footer)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
